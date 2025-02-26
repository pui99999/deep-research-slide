from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import datetime
import os

# カラースキームの定義 - よりモダンな配色に更新
TEAL_BLUE = RGBColor(0, 150, 199)       # #0096C7 (主色) - より鮮やかなブルー
LIGHT_GRAY = RGBColor(245, 247, 249)    # #F5F7F9 (補色)
LEAF_GREEN = RGBColor(80, 184, 72)      # #50B848 (アクセント1) - より鮮やかなグリーン
GRAPHITE = RGBColor(66, 66, 66)         # #424242 (アクセント2)
AMBER = RGBColor(255, 149, 0)           # #FF9500 (強調色) - よりモダンなオレンジ
DARK_GRAY = RGBColor(51, 51, 51)        # #333333 (本文)
WHITE = RGBColor(255, 255, 255)         # #FFFFFF (白)
LIGHT_BLUE = RGBColor(230, 246, 255)    # #E6F6FF (背景色) - 新しい色

# タイポグラフィサイズの定義
TITLE_SIZE = Pt(40)
HEADING_SIZE = Pt(32)
SUB_HEADING_SIZE = Pt(24)
BODY_SIZE = Pt(20)
SMALL_SIZE = Pt(18)

def apply_slide_background(slide):
    """すべてのスライドに共通の背景とデザイン要素を適用"""
    # 背景色を設定
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # 左側のアクセントバー
    left = Inches(0)
    top = Inches(0)
    width = Inches(0.5)
    height = Inches(7.5)
    rect = slide.shapes.add_shape(
        1, left, top, width, height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = TEAL_BLUE
    rect.line.fill.background()
    
    # 右上の装飾円
    left = Inches(9)
    top = Inches(0.2)
    width = Inches(0.8)
    height = Inches(0.8)
    oval = slide.shapes.add_shape(
        3, left, top, width, height
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = LIGHT_BLUE
    oval.line.fill.background()
    
    # 右下の装飾円
    left = Inches(9.2)
    top = Inches(6.5)
    width = Inches(0.6)
    height = Inches(0.6)
    oval = slide.shapes.add_shape(
        3, left, top, width, height
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = LEAF_GREEN
    oval.line.fill.background()

def setup_title_slide(slide):
    """タイトルスライドのセットアップ"""
    # 背景色を設定
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BLUE
    
    # 装飾要素 - 左側の縦線
    left = Inches(1.2)
    top = Inches(1.5)
    width = Inches(0.1)
    height = Inches(4.5)
    rect = slide.shapes.add_shape(
        1, left, top, width, height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = TEAL_BLUE
    rect.line.fill.background()
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    # タイトル設定
    title_shape.text = "日本の産業廃棄物処理業界の\n市場規模に関する調査"
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.LEFT
    title_para.font.size = TITLE_SIZE
    title_para.font.color.rgb = TEAL_BLUE
    title_para.font.name = 'Noto Sans JP'
    title_para.font.bold = True
    
    # サブタイトル（日付）設定
    today = datetime.datetime.now().strftime("%Y年%m月")
    subtitle_shape.text = today
    subtitle_para = subtitle_shape.text_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.LEFT
    subtitle_para.font.size = SMALL_SIZE
    subtitle_para.font.color.rgb = GRAPHITE
    subtitle_para.font.name = 'Noto Sans JP'
    
    # リサイクルアイコン（テキストで代用）
    left = Inches(8)
    top = Inches(5.5)
    width = Inches(1)
    height = Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = "♻"  # リサイクルアイコン
    p.font.size = Pt(60)
    p.font.color.rgb = TEAL_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 装飾要素 - 右上の円形
    left = Inches(8)
    top = Inches(0.5)
    width = Inches(1.5)
    height = Inches(1.5)
    oval = slide.shapes.add_shape(
        3, left, top, width, height
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = LEAF_GREEN
    oval.line.fill.background()
    
    # 装飾要素 - 左下の円形
    left = Inches(0.5)
    top = Inches(6)
    width = Inches(0.8)
    height = Inches(0.8)
    oval = slide.shapes.add_shape(
        3, left, top, width, height
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = AMBER
    oval.line.fill.background()
    
    # 区切り線
    left = Inches(1.5)
    top = Inches(4.5)
    width = Inches(4)
    height = Inches(0.05)
    line = slide.shapes.add_shape(
        1, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LEAF_GREEN
    line.line.fill.background()

def add_overview_slide(prs):
    """調査概要スライドの追加"""
    slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツのレイアウト
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_background(slide)
    
    # タイトル設定
    title_shape = slide.shapes.title
    title_shape.text = "調査概要"
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.font.size = HEADING_SIZE
    title_para.font.color.rgb = TEAL_BLUE
    title_para.font.name = 'Noto Sans JP'
    title_para.font.bold = True
    
    # 細い水平線
    left = Inches(1)
    top = Inches(1.3)
    width = Inches(8)
    height = Inches(0.02)
    line = slide.shapes.add_shape(
        1, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL_BLUE
    line.line.fill.background()
    
    # コンテンツ - Y位置を上に調整
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "本調査では以下の情報を含みます:"
    tf.paragraphs[0].font.size = BODY_SIZE
    tf.paragraphs[0].font.color.rgb = DARK_GRAY
    tf.paragraphs[0].font.name = 'Noto Sans JP'
    
    items = [
        "最新の市場規模（売上・成長率）", 
        "産業別の廃棄物排出量", 
        "処理方法別の市場規模", 
        "主要な事業者とシェア", 
        "政府の規制や補助金の影響"
    ]
    
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = BODY_SIZE
        p.font.color.rgb = GRAPHITE
        p.font.name = 'Noto Sans JP'
        p.level = 1
        
    # 矢印アイコン（テキストで代用）- 位置調整
    left = Inches(8)
    top = Inches(5.5)  # 6から5.5に上方向に調整
    width = Inches(1)
    height = Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = "↓"  # 矢印アイコン
    p.font.size = Pt(36)
    p.font.color.rgb = TEAL_BLUE
    p.alignment = PP_ALIGN.CENTER

def add_market_size_slide(prs):
    """市場規模スライドの追加"""
    slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツのレイアウト
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_background(slide)
    
    # タイトル設定
    title_shape = slide.shapes.title
    title_shape.text = "市場規模（最新の動向）"
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.font.size = HEADING_SIZE
    title_para.font.color.rgb = TEAL_BLUE
    title_para.font.name = 'Noto Sans JP'
    title_para.font.bold = True
    
    # 細い水平線
    left = Inches(1)
    top = Inches(1.3)
    width = Inches(8)
    height = Inches(0.02)
    line = slide.shapes.add_shape(
        1, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL_BLUE
    line.line.fill.background()
    
    # 折れ線グラフのデータ
    chart_data = CategoryChartData()
    chart_data.categories = ['2013', '2014', '2015', '2016', '2017', '2018', '2019', '2020', '2021']
    chart_data.add_series('売上高（兆円）', (1.8, 1.9, 2.0, 2.1, 2.2, 2.4, 2.5, 2.66, 2.8))
    
    # グラフの追加 - 位置調整
    x, y, cx, cy = Inches(1.5), Inches(1.8), Inches(7), Inches(3)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart
    
    # グラフのスタイル設定
    line_series = chart.series[0]
    line_series.format.line.color.rgb = TEAL_BLUE
    line_series.format.line.width = Pt(3)
    
    # グラフの枠線を削除
    chart.has_border = False
    
    # 箇条書きテキスト - 位置調整
    left = Inches(1.5)
    top = Inches(5)
    width = Inches(7)
    height = Inches(1.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    
    bullet_items = [
        "2020年: 2兆6,634億円（前年比+6.7%）",
        "2013年から一貫した増加傾向",
        "2020年は初の2.6兆円台突破"
    ]
    
    for i, item in enumerate(bullet_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = "• " + item
        p.font.size = BODY_SIZE
        p.font.color.rgb = DARK_GRAY
        p.font.name = 'Noto Sans JP'
        
        # 強調したい数字をハイライト
        if "2.6兆円" in item:
            run = p.add_run()
            run.text = " (初)"
            run.font.color.rgb = AMBER
            run.font.bold = True
            
    # 装飾要素 - 右側の縦線
    left = Inches(8.8)
    top = Inches(2)
    width = Inches(0.05)
    height = Inches(3.5)
    line = slide.shapes.add_shape(
        1, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LEAF_GREEN
    line.line.fill.background()

def add_industry_breakdown_slide(prs):
    """産業別の廃棄物排出量スライドの追加"""
    slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツのレイアウト
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_background(slide)
    
    # タイトル設定
    title_shape = slide.shapes.title
    title_shape.text = "産業別の廃棄物排出量"
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.font.size = HEADING_SIZE
    title_para.font.color.rgb = TEAL_BLUE
    title_para.font.name = 'Noto Sans JP'
    title_para.font.bold = True
    
    # 細い水平線
    left = Inches(1)
    top = Inches(1.3)
    width = Inches(8)
    height = Inches(0.02)
    line = slide.shapes.add_shape(
        1, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL_BLUE
    line.line.fill.background()
    
    # 円グラフのデータ
    chart_data = CategoryChartData()
    chart_data.categories = [
        '電気・ガス・水道', '農業・林業', '建設業', 
        'パルプ・紙工業', '鉄鋼業', 'その他'
    ]
    chart_data.add_series('排出割合', (26.5, 21.7, 21.5, 7.4, 6.2, 16.7))
    
    # グラフの追加 - 位置調整
    x, y, cx, cy = Inches(1.2), Inches(1.8), Inches(4), Inches(3.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart
    
    # グラフの枠線を削除
    chart.has_border = False
    
    # 円グラフのスライスの色を設定
    slices = chart.plots[0].series[0].points
    slice_colors = [TEAL_BLUE, LEAF_GREEN, AMBER, RGBColor(100, 181, 246), RGBColor(121, 85, 72), GRAPHITE]
    for i, slice in enumerate(slices):
        slice.format.fill.solid()
        slice.format.fill.fore_color.rgb = slice_colors[i % len(slice_colors)]
    
    # サブタイトル - 位置調整
    left = Inches(5.5)
    top = Inches(2.0)
    width = Inches(4)
    height = Inches(0.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = "上位5業種（総排出量の83%）:"
    tf.paragraphs[0].font.size = SUB_HEADING_SIZE
    tf.paragraphs[0].font.color.rgb = GRAPHITE
    tf.paragraphs[0].font.name = 'Noto Sans JP'
    tf.paragraphs[0].font.bold = True
    
    # 箇条書きテキスト - 位置調整
    left = Inches(5.5)
    top = Inches(2.6)
    width = Inches(4)
    height = Inches(3)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    
    # 産業とパーセンテージのリスト
    industries = [
        "電気・ガス・水道: 26.5%",
        "農業・林業: 21.7%",
        "建設業: 21.5%",
        "パルプ・紙工業: 7.4%",
        "鉄鋼業: 6.2%"
    ]
    
    for i, industry in enumerate(industries):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = "• " + industry
        p.font.size = BODY_SIZE
        p.font.color.rgb = DARK_GRAY
        p.font.name = 'Noto Sans JP'
        
        # 行間隔を調整
        p.space_after = Pt(10)

def add_treatment_methods_slide(prs):
    """処理方法別の内訳スライドの追加"""
    slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツのレイアウト
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_background(slide)
    
    # タイトル設定
    title_shape = slide.shapes.title
    title_shape.text = "処理方法別の内訳"
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.font.size = HEADING_SIZE
    title_para.font.color.rgb = TEAL_BLUE
    title_para.font.name = 'Noto Sans JP'
    title_para.font.bold = True
    
    # 細い水平線
    left = Inches(1)
    top = Inches(1.3)
    width = Inches(8)
    height = Inches(0.02)
    line = slide.shapes.add_shape(
        1, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL_BLUE
    line.line.fill.background()
    
    # 棒グラフのデータ
    chart_data = CategoryChartData()
    chart_data.categories = ['再生利用', '焼却等の中間処理', '最終処分（埋立）']
    chart_data.add_series('処理割合 (%)', (54.2, 43.5, 2.3))
    
    # グラフの追加 - 位置を上に調整してサイズも小さく
    x, y, cx, cy = Inches(1.0), Inches(1.7), Inches(4.5), Inches(2.3)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    # グラフのスタイル設定
    bar_series = chart.series[0]
    bar_fill = bar_series.format.fill
    bar_fill.solid()
    bar_fill.fore_color.rgb = TEAL_BLUE
    
    # グラフの枠線を削除
    chart.has_border = False
    
    # 処理方法とデータの詳細テキスト
    treatment_details = [
        ("再生利用（リサイクル）: 54.2%", "(2億0372万トン)"),
        ("焼却等の中間処理: 43.5%", "(1億6337万トン)"),
        ("最終処分（埋立）: 2.3%", "(883万トン)")
    ]
    
    # テキストボックスの位置調整
    left = Inches(5.8)
    top = Inches(1.7)
    width = Inches(3.5)
    height = Inches(3.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    
    # 最初の段落を作成
    p = tf.paragraphs[0]
    p.text = ""  # 一旦空にする
    
    for i, (method, amount) in enumerate(treatment_details):
        if i > 0:
            p = tf.add_paragraph()
        
        # 処理方法と割合を強調
        run = p.add_run()
        run.text = "• " + method
        run.font.size = BODY_SIZE
        run.font.color.rgb = DARK_GRAY
        run.font.name = 'Noto Sans JP'
        run.font.bold = True
        
        # トン数を通常のスタイルで、行間調整
        p.add_line_break()
        run = p.add_run()
        run.text = "   " + amount  # インデント用のスペース
        run.font.size = BODY_SIZE
        run.font.color.rgb = DARK_GRAY
        run.font.name = 'Noto Sans JP'
        run.font.bold = False
        
        # 行間隔を調整
        p.space_after = Pt(15)
    
    # 装飾要素 - 下部の横線
    left = Inches(1.0)
    top = Inches(4.5)
    width = Inches(8)
    height = Inches(0.02)
    line = slide.shapes.add_shape(
        1, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LEAF_GREEN
    line.line.fill.background()
    
    # リサイクルアイコン
    left = Inches(4.0)
    top = Inches(5.0)
    width = Inches(1)
    height = Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = "♻"  # リサイクルアイコン
    p.font.size = Pt(60)
    p.font.color.rgb = TEAL_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 説明テキスト
    left = Inches(2.0)
    top = Inches(6.0)
    width = Inches(6)
    height = Inches(0.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = "リサイクル率は年々向上し、最終処分量は大幅に減少しています"
    p.font.size = SMALL_SIZE
    p.font.color.rgb = GRAPHITE
    p.font.name = 'Noto Sans JP'
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

def add_major_companies_slide(prs):
    """主な事業者と市場シェアスライドの追加"""
    slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツのレイアウト
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_background(slide)
    
    # タイトル設定
    title_shape = slide.shapes.title
    title_shape.text = "主な事業者と市場シェア"
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.font.size = HEADING_SIZE
    title_para.font.color.rgb = TEAL_BLUE
    title_para.font.name = 'Noto Sans JP'
    title_para.font.bold = True
    
    # 細い水平線
    left = Inches(1)
    top = Inches(1.3)
    width = Inches(8)
    height = Inches(0.02)
    line = slide.shapes.add_shape(
        1, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL_BLUE
    line.line.fill.background()
    
    # ドーナツチャートのデータ（簡易的に円グラフで代用）
    chart_data = CategoryChartData()
    chart_data.categories = ['TREホールディングス', 'エンビプロHD', 'ダイセキ', 'その他主要3社', '残りの業者']
    chart_data.add_series('シェア (%)', (2, 1.5, 1.5, 2, 93))
    
    # グラフの追加 - サイズと位置調整
    x, y, cx, cy = Inches(1.0), Inches(1.7), Inches(3.5), Inches(3.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
    ).chart
    
    # グラフの枠線を削除
    chart.has_border = False
    
    # ドーナツチャートのスライスの色を設定
    slices = chart.plots[0].series[0].points
    slice_colors = [TEAL_BLUE, LEAF_GREEN, AMBER, RGBColor(100, 181, 246), LIGHT_GRAY]
    for i, slice in enumerate(slices):
        slice.format.fill.solid()
        slice.format.fill.fore_color.rgb = slice_colors[i % len(slice_colors)]
    
    # 右側のコンテンツエリア
    # サブタイトル: 業界特性 - 位置調整
    left = Inches(5.0)
    top = Inches(1.7)
    width = Inches(4.5)
    height = Inches(0.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = "業界特性:"
    tf.paragraphs[0].font.size = SUB_HEADING_SIZE
    tf.paragraphs[0].font.color.rgb = GRAPHITE
    tf.paragraphs[0].font.name = 'Noto Sans JP'
    tf.paragraphs[0].font.bold = True
    
    # 業界特性の箇条書き - 位置調整
    left = Inches(5.0)
    top = Inches(2.2)
    width = Inches(4.5)
    height = Inches(1.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    
    characteristics = [
        "全国で約12万社の事業者",
        "上位企業でも市場の一部を占めるのみ",
        "主要6社の合計シェア: 約7%"
    ]
    
    for i, item in enumerate(characteristics):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = "• " + item
        p.font.size = BODY_SIZE
        p.font.color.rgb = DARK_GRAY
        p.font.name = 'Noto Sans JP'
        p.space_after = Pt(10)  # 行間隔を調整
        
        # 強調すべき数字
        if "12万社" in item or "7%" in item:
            parts = p.text.split(": " if ": " in item else " ")
            p.text = parts[0]
            
            # 強調テキストを追加
            if ": " in item:
                p.add_run().text = ": "
                run = p.add_run()
                run.text = parts[1]
                run.font.color.rgb = AMBER
                run.font.bold = True
            else:
                for j, part in enumerate(parts[1:]):
                    if j > 0 or "社" not in part:
                        p.add_run().text = " "
                    run = p.add_run()
                    run.text = part
                    if "12万" in part or "7%" in part:
                        run.font.color.rgb = AMBER
                        run.font.bold = True
                    else:
                        run.font.color.rgb = DARK_GRAY
    
    # 売上上位企業 - 位置調整
    left = Inches(5.0)
    top = Inches(4.0)
    width = Inches(4.5)
    height = Inches(0.4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = "売上上位企業（2021年）:"
    tf.paragraphs[0].font.size = SUB_HEADING_SIZE
    tf.paragraphs[0].font.color.rgb = GRAPHITE
    tf.paragraphs[0].font.name = 'Noto Sans JP'
    tf.paragraphs[0].font.bold = True
    
    # 上位企業リスト - 位置調整
    left = Inches(5.0)
    top = Inches(4.5)
    width = Inches(4.5)
    height = Inches(1.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    
    companies = [
        "TREホールディングス: 682億円",
        "エンビプロHD: 573億円",
        "ダイセキ: 568億円"
    ]
    
    for i, company in enumerate(companies):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = "• " + company
        p.font.size = BODY_SIZE
        p.font.color.rgb = DARK_GRAY
        p.font.name = 'Noto Sans JP'
        p.space_after = Pt(10)  # 行間隔を調整
    
    # 装飾要素 - 下部の横線
    left = Inches(1.0)
    top = Inches(6.0)
    width = Inches(8)
    height = Inches(0.02)
    line = slide.shapes.add_shape(
        1, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LEAF_GREEN
    line.line.fill.background()
    
    # 説明テキスト
    left = Inches(1.0)
    top = Inches(6.2)
    width = Inches(8)
    height = Inches(0.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = "業界は非常に分散しており、地域密着型の中小企業が多数存在"
    p.font.size = SMALL_SIZE
    p.font.color.rgb = GRAPHITE
    p.font.name = 'Noto Sans JP'
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

def add_government_regulations_slide(prs):
    """政府の規制と業界への影響スライドの追加"""
    slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツのレイアウト
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_background(slide)
    
    # タイトル設定
    title_shape = slide.shapes.title
    title_shape.text = "政府の規制と業界への影響"
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.font.size = HEADING_SIZE
    title_para.font.color.rgb = TEAL_BLUE
    title_para.font.name = 'Noto Sans JP'
    title_para.font.bold = True
    
    # 細い水平線
    left = Inches(1)
    top = Inches(1.3)
    width = Inches(8)
    height = Inches(0.02)
    line = slide.shapes.add_shape(
        1, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL_BLUE
    line.line.fill.background()
    
    # 左側: 規制の影響 - 位置調整
    left = Inches(1.0)
    top = Inches(1.7)
    width = Inches(4.0)
    height = Inches(0.4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = "規制の影響:"
    tf.paragraphs[0].font.size = SUB_HEADING_SIZE
    tf.paragraphs[0].font.color.rgb = GRAPHITE
    tf.paragraphs[0].font.name = 'Noto Sans JP'
    tf.paragraphs[0].font.bold = True
    
    # 規制の影響リスト
    regulations = [
        "廃棄物処理法",
        "マニフェスト制度",
        "最終処分量85%減（1997→2014年）",
        "不法投棄件数減少",
        "リサイクル率向上(50%前後で安定)"
    ]
    
    # 位置とサイズ調整
    left = Inches(1.0)
    top = Inches(2.2)
    width = Inches(4.0)
    height = Inches(3)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    
    for i, item in enumerate(regulations):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = "• " + item
        p.font.size = BODY_SIZE
        p.font.color.rgb = DARK_GRAY
        p.font.name = 'Noto Sans JP'
        p.space_after = Pt(10)  # 行間隔を調整
        
        # 強調すべき数字
        if "85%" in item:
            parts = item.split("85%")
            p.text = "• " + parts[0]
            run = p.add_run()
            run.text = "85%"
            run.font.color.rgb = AMBER
            run.font.bold = True
            p.add_run().text = parts[1]
    
    # 中央の区切り線 - 位置調整
    left = Inches(5.0)
    top = Inches(1.7)
    width = Inches(0.02)
    height = Inches(4)
    line = slide.shapes.add_shape(
        1, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL_BLUE
    line.line.fill.background()
    
    # 右側: 補助金制度 - 位置調整
    left = Inches(5.5)
    top = Inches(1.7)
    width = Inches(4.0)
    height = Inches(0.4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = "補助金制度:"
    tf.paragraphs[0].font.size = SUB_HEADING_SIZE
    tf.paragraphs[0].font.color.rgb = GRAPHITE
    tf.paragraphs[0].font.name = 'Noto Sans JP'
    tf.paragraphs[0].font.bold = True
    
    # 補助金制度リスト
    subsidies = [
        "産業廃棄物処理事業振興財団",
        "技術開発補助金",
        "設備投資支援"
    ]
    
    # 位置とサイズ調整
    left = Inches(5.5)
    top = Inches(2.2)
    width = Inches(4.0)
    height = Inches(3)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    
    for i, item in enumerate(subsidies):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = "• " + item
        p.font.size = BODY_SIZE
        p.font.color.rgb = DARK_GRAY
        p.font.name = 'Noto Sans JP'
        p.space_after = Pt(10)  # 行間隔を調整
    
    # 装飾要素 - 下部の横線
    left = Inches(1.0)
    top = Inches(6.0)
    width = Inches(8)
    height = Inches(0.02)
    line = slide.shapes.add_shape(
        1, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LEAF_GREEN
    line.line.fill.background()
    
    # 説明テキスト
    left = Inches(1.0)
    top = Inches(6.2)
    width = Inches(8)
    height = Inches(0.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = "政府の規制と支援が業界の健全な発展を促進しています"
    p.font.size = SMALL_SIZE
    p.font.color.rgb = GRAPHITE
    p.font.name = 'Noto Sans JP'
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

def add_summary_slide(prs):
    """まとめスライドの追加"""
    slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツのレイアウト
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_background(slide)
    
    # タイトル設定
    title_shape = slide.shapes.title
    title_shape.text = "まとめ"
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.font.size = HEADING_SIZE
    title_para.font.color.rgb = TEAL_BLUE
    title_para.font.name = 'Noto Sans JP'
    title_para.font.bold = True
    
    # 細い水平線
    left = Inches(1)
    top = Inches(1.3)
    width = Inches(8)
    height = Inches(0.02)
    line = slide.shapes.add_shape(
        1, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL_BLUE
    line.line.fill.background()
    
    # まとめの箇条書き
    summary_items = [
        "産業廃棄物処理は5兆円規模の大きな市場",
        "上位5業種で排出量の8割以上",
        "リサイクル率は54%超、最終処分はわずか2.3%",
        "分散型市場構造（約12万社）",
        "政府規制が業界発展を下支え"
    ]
    
    # 左側の装飾要素 - 縦線
    left = Inches(1.0)
    top = Inches(1.7)
    width = Inches(0.05)
    height = Inches(4.5)
    line = slide.shapes.add_shape(
        1, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LEAF_GREEN
    line.line.fill.background()
    
    # 右側の装飾要素 - 縦線
    left = Inches(8.5)
    top = Inches(1.7)
    width = Inches(0.05)
    height = Inches(4.5)
    line = slide.shapes.add_shape(
        1, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LEAF_GREEN
    line.line.fill.background()
    
    # 位置とサイズ調整 - 上部スペースを増やす
    left = Inches(1.5)
    top = Inches(2.0)
    width = Inches(7)
    height = Inches(3.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True  # 単語の折り返しを有効化
    
    for i, item in enumerate(summary_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            # 段落間のスペースを追加
            if i > 0:
                p.space_before = Pt(15)
        
        # 強調したいキーワードを判断してスタイルを変える
        if "5兆円" in item:
            parts = item.split("5兆円")
            p.text = "• " + parts[0]
            run = p.add_run()
            run.text = "5兆円"
            run.font.color.rgb = AMBER
            run.font.bold = True
            p.add_run().text = parts[1]
        elif "8割" in item:
            parts = item.split("8割")
            p.text = "• " + parts[0]
            run = p.add_run()
            run.text = "8割"
            run.font.color.rgb = AMBER
            run.font.bold = True
            p.add_run().text = parts[1]
        elif "54%" in item or "2.3%" in item:
            p.text = "• " + item.split("54%")[0]
            run = p.add_run()
            run.text = "54%"
            run.font.color.rgb = AMBER
            run.font.bold = True
            middle_text = item.split("54%")[1].split("2.3%")[0]
            p.add_run().text = middle_text
            run = p.add_run()
            run.text = "2.3%"
            run.font.color.rgb = AMBER
            run.font.bold = True
            p.add_run().text = item.split("2.3%")[1]
        elif "12万社" in item:
            parts = item.split("12万社")
            p.text = "• " + parts[0] + "（"
            run = p.add_run()
            run.text = "約12万社"
            run.font.color.rgb = AMBER
            run.font.bold = True
            p.add_run().text = "）"
        else:
            p.text = "• " + item
        
        p.font.size = SUB_HEADING_SIZE
        p.font.color.rgb = GRAPHITE
        p.font.name = 'Noto Sans JP'
        p.font.bold = True
    
    # リサイクルアイコン（テキストで代用）- 位置調整
    left = Inches(4.5)
    top = Inches(6.0)
    width = Inches(1)
    height = Inches(1)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = "♻"  # リサイクルアイコン
    p.font.size = Pt(60)
    p.font.color.rgb = TEAL_BLUE
    p.alignment = PP_ALIGN.CENTER

def create_presentation():
    """プレゼンテーションの作成"""
    prs = Presentation()
    
    # スライドサイズをワイドスクリーンに設定
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # タイトルスライドの追加
    title_slide_layout = prs.slide_layouts[0]  # タイトルスライドのレイアウト
    title_slide = prs.slides.add_slide(title_slide_layout)
    setup_title_slide(title_slide)
    
    # 各セクションのスライドを追加
    add_overview_slide(prs)
    add_market_size_slide(prs)
    add_industry_breakdown_slide(prs)
    add_treatment_methods_slide(prs)
    add_major_companies_slide(prs)
    add_government_regulations_slide(prs)
    add_summary_slide(prs)
    
    # プレゼンテーションを保存
    today_str = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"産業廃棄物市場分析_{today_str}.pptx"
    prs.save(filename)
    print(f"プレゼンテーションを保存しました: {filename}")
    
    return filename

if __name__ == "__main__":
    output_file = create_presentation()
    print(f"プレゼンテーションが正常に作成されました。ファイル: {output_file}") 