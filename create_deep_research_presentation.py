from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re
import os
import argparse
from datetime import datetime

def create_deep_research_presentation(research_text, output_file=None, title="研究結果", theme="blue"):
    """
    OpenAIのDeepResearchの結果をパワーポイントにまとめる関数
    
    Args:
        research_text (str): DeepResearchの結果テキスト
        output_file (str): 出力するパワーポイントファイル名（Noneの場合は自動生成）
        title (str): プレゼンテーションのタイトル
        theme (str): カラーテーマ（"blue", "dark", "light", "green"）
    """
    # 出力ファイル名が指定されていない場合は自動生成
    if output_file is None:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_file = f"deep_research_{timestamp}.pptx"
    
    # プレゼンテーションの作成
    prs = Presentation()
    
    # スライドのレイアウト
    title_slide_layout = prs.slide_layouts[0]  # タイトルスライド
    content_slide_layout = prs.slide_layouts[1]  # タイトルと内容のスライド
    
    # カラーテーマの設定
    if theme == "blue":
        title_color = RGBColor(0, 112, 192)  # 青
        subtitle_color = RGBColor(0, 176, 240)  # 明るい青
        text_color = RGBColor(0, 0, 0)  # 黒
        highlight_color = RGBColor(192, 0, 0)  # 赤
        background_color = RGBColor(240, 240, 240)  # 薄いグレー
    elif theme == "dark":
        title_color = RGBColor(255, 255, 255)  # 白
        subtitle_color = RGBColor(200, 200, 200)  # 薄いグレー
        text_color = RGBColor(255, 255, 255)  # 白
        highlight_color = RGBColor(255, 128, 0)  # オレンジ
        background_color = RGBColor(44, 44, 44)  # 暗いグレー
    elif theme == "light":
        title_color = RGBColor(70, 70, 70)  # 暗いグレー
        subtitle_color = RGBColor(100, 100, 100)  # グレー
        text_color = RGBColor(0, 0, 0)  # 黒
        highlight_color = RGBColor(255, 128, 0)  # オレンジ
        background_color = RGBColor(255, 255, 255)  # 白
    elif theme == "green":
        title_color = RGBColor(0, 128, 0)  # 緑
        subtitle_color = RGBColor(0, 176, 80)  # 明るい緑
        text_color = RGBColor(0, 0, 0)  # 黒
        highlight_color = RGBColor(192, 0, 0)  # 赤
        background_color = RGBColor(240, 240, 240)  # 薄いグレー
    else:
        # デフォルトは青テーマ
        title_color = RGBColor(0, 112, 192)  # 青
        subtitle_color = RGBColor(0, 176, 240)  # 明るい青
        text_color = RGBColor(0, 0, 0)  # 黒
        highlight_color = RGBColor(192, 0, 0)  # 赤
        background_color = RGBColor(240, 240, 240)  # 薄いグレー
    
    # 背景画像のパス（オプション）
    bg_image_path = None
    
    # タイトルスライドの作成
    title_slide = prs.slides.add_slide(title_slide_layout)
    
    # 背景色の設定
    if theme == "dark":
        # 暗いテーマの場合は背景を暗く
        background = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color
        background.line.fill.background()
        # 背景を最背面に配置
        background.shadow.inherit = False
        background.zorder = 0
    
    # タイトルとサブタイトルのテキストボックス
    title_shape = title_slide.shapes.title
    subtitle_shape = title_slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = "研究結果プレゼンテーション"
    
    # タイトルの色を設定
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = title_color
            run.font.size = Pt(44)
            run.font.bold = True
    
    # サブタイトルの色を設定
    for paragraph in subtitle_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = subtitle_color
            run.font.size = Pt(28)
            run.font.italic = True
    
    # 日付を追加
    date_box = title_slide.shapes.add_textbox(
        Inches(0.5), Inches(5), Inches(9), Inches(0.5)
    )
    date_tf = date_box.text_frame
    date_p = date_tf.add_paragraph()
    date_p.text = datetime.now().strftime("%Y年%m月%d日")
    date_p.alignment = PP_ALIGN.RIGHT
    date_run = date_p.runs[0]
    date_run.font.size = Pt(12)
    date_run.font.color.rgb = subtitle_color
    
    # URLを除外するための正規表現パターン
    url_pattern = re.compile(r'https?://\S+|www\.\S+|\[\d+\]|\(\d+\)|参考文献|References')
    
    # 参考文献セクションを検出するパターン
    ref_section_pattern = re.compile(r'^参考文献|^References|^引用文献|^Sources|^Citations', re.IGNORECASE)
    
    # テキストを段落に分割
    paragraphs = re.split(r'\n\s*\n', research_text)
    
    # 参考文献セクションを除外
    filtered_paragraphs = []
    for para in paragraphs:
        # 参考文献セクションを検出したら、それ以降は含めない
        if any(ref_section_pattern.match(line.strip()) for line in para.split('\n')):
            break
        filtered_paragraphs.append(para)
    
    paragraphs = filtered_paragraphs
    
    # 目次スライドの作成
    toc_slide = prs.slides.add_slide(content_slide_layout)
    
    # 背景色の設定
    if theme == "dark":
        # 暗いテーマの場合は背景を暗く
        background = toc_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color
        background.line.fill.background()
        # 背景を最背面に配置
        background.shadow.inherit = False
        background.zorder = 0
    
    toc_title = toc_slide.shapes.title
    toc_content = toc_slide.placeholders[1]
    
    toc_title.text = "目次"
    
    # タイトルの色を設定
    for paragraph in toc_title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = title_color
            run.font.size = Pt(40)
            run.font.bold = True
    
    # 目次の内容を作成
    toc_text = ""
    section_titles = []
    
    for i, para in enumerate(paragraphs):
        if i == 0:  # 最初の段落はタイトルなのでスキップ
            continue
        
        # 段落の最初の行をセクションタイトルとして使用
        lines = para.strip().split('\n')
        if lines:
            section_title = lines[0].strip()
            # URLや参考文献を含む行は除外
            if not url_pattern.search(section_title):
                # 数字だけの見出しや短すぎる見出しは除外
                if not re.match(r'^\d+\.?$', section_title) and len(section_title) > 3:
                    section_titles.append(section_title)
                    toc_text += f"• {section_title}\n"
    
    toc_content.text = toc_text
    
    # 目次テキストの色を設定
    for paragraph in toc_content.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = text_color
            run.font.size = Pt(24)
    
    # 各セクションのスライドを作成
    current_section = ""
    section_content = ""
    
    for i, para in enumerate(paragraphs):
        if i == 0:  # 最初の段落はタイトルなのでスキップ
            continue
        
        lines = para.strip().split('\n')
        if not lines:
            continue
        
        # URLや参考文献を含む行は除外
        filtered_lines = []
        for line in lines:
            # URLや引用番号を含まない行だけを追加
            if not url_pattern.search(line):
                filtered_lines.append(line)
        
        if not filtered_lines:
            continue
        
        section_title = filtered_lines[0].strip()
        
        # 新しいセクションの開始
        if section_title in section_titles:
            # 前のセクションがあれば、そのスライドを作成
            if current_section:
                create_research_slide(prs, content_slide_layout, current_section, section_content, 
                                     title_color, text_color, background_color, theme)
            
            current_section = section_title
            section_content = '\n'.join(filtered_lines[1:])
        else:
            # 同じセクションの続き
            section_content += '\n' + '\n'.join(filtered_lines)
    
    # 最後のセクションのスライドを作成
    if current_section:
        create_research_slide(prs, content_slide_layout, current_section, section_content, 
                             title_color, text_color, background_color, theme)
    
    # まとめスライドの作成
    summary_slide = prs.slides.add_slide(content_slide_layout)
    
    # 背景色の設定
    if theme == "dark":
        # 暗いテーマの場合は背景を暗く
        background = summary_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color
        background.line.fill.background()
        # 背景を最背面に配置
        background.shadow.inherit = False
        background.zorder = 0
    
    summary_title = summary_slide.shapes.title
    summary_content = summary_slide.placeholders[1]
    
    summary_title.text = "まとめ"
    
    # タイトルの色を設定
    for paragraph in summary_title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = title_color
            run.font.size = Pt(40)
            run.font.bold = True
    
    # まとめの内容（最初の段落から抽出するか、固定テキスト）
    if len(paragraphs) > 0:
        # 最初の段落から要約を抽出（URLや参考文献を除外）
        first_para = paragraphs[0]
        summary_lines = [line for line in first_para.split('\n') if not url_pattern.search(line)]
        summary_text = '\n'.join(summary_lines)
        
        # 長すぎる場合は短縮
        if len(summary_text) > 500:
            summary_text = summary_text[:497] + "..."
    else:
        summary_text = "研究結果の主要ポイント"
    
    summary_content.text = summary_text
    
    # まとめテキストの色を設定
    for paragraph in summary_content.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = text_color
            run.font.size = Pt(24)
    
    # プレゼンテーションの保存
    prs.save(output_file)
    print(f"研究プレゼンテーションを {output_file} として保存しました。")
    
    return output_file

def create_research_slide(prs, layout, title, content, title_color, text_color, background_color, theme):
    """
    研究結果のスライドを作成する関数
    
    Args:
        prs: プレゼンテーションオブジェクト
        layout: スライドレイアウト
        title (str): スライドのタイトル
        content (str): スライドの内容
        title_color: タイトルの色
        text_color: テキストの色
        background_color: 背景色
        theme: カラーテーマ
    """
    slide = prs.slides.add_slide(layout)
    
    # 背景色の設定
    if theme == "dark":
        # 暗いテーマの場合は背景を暗く
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color
        background.line.fill.background()
        # 背景を最背面に配置
        background.shadow.inherit = False
        background.zorder = 0
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    
    # タイトルの色を設定
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = title_color
            run.font.size = Pt(36)
            run.font.bold = True
    
    # 内容が長すぎる場合は分割
    if len(content) > 1500:
        # 複数のスライドに分割する必要がある
        lines = content.split('\n')
        chunks = []
        current_chunk = []
        current_length = 0
        
        for line in lines:
            line_length = len(line)
            if current_length + line_length > 1500:
                chunks.append('\n'.join(current_chunk))
                current_chunk = [line]
                current_length = line_length
            else:
                current_chunk.append(line)
                current_length += line_length
        
        if current_chunk:
            chunks.append('\n'.join(current_chunk))
        
        # 最初のチャンクをこのスライドに表示
        content = chunks[0]
        
        # 残りのチャンクを新しいスライドに表示
        for i, chunk in enumerate(chunks[1:], 1):
            continuation_slide = prs.slides.add_slide(layout)
            
            # 背景色の設定
            if theme == "dark":
                background = continuation_slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
                )
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = background_color
                background.line.fill.background()
                # 背景を最背面に配置
                background.shadow.inherit = False
                background.zorder = 0
            
            cont_title = continuation_slide.shapes.title
            cont_content = continuation_slide.placeholders[1]
            
            cont_title.text = f"{title} (続き {i})"
            
            # タイトルの色を設定
            for paragraph in cont_title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = title_color
                    run.font.size = Pt(36)
                    run.font.bold = True
            
            cont_content.text = chunk
            
            # テキストの色を設定
            for paragraph in cont_content.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = text_color
                    run.font.size = Pt(18)
    
    content_shape.text = content
    
    # テキストの色を設定
    for paragraph in content_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = text_color
            run.font.size = Pt(18)
    
    return slide

def extract_title_from_text(text):
    """
    テキストから適切なタイトルを抽出する
    
    Args:
        text (str): 入力テキスト
        
    Returns:
        str: 抽出されたタイトル
    """
    # 最初の行または段落をタイトルとして使用
    lines = text.strip().split('\n')
    if lines:
        first_line = lines[0].strip()
        # URLや参考文献を含まない場合のみ使用
        if not re.search(r'https?://|www\.|\[\d+\]|\(\d+\)|参考文献|References', first_line):
            return first_line
    
    # 適切なタイトルが見つからない場合はデフォルト
    return "研究結果"

def main():
    parser = argparse.ArgumentParser(description='OpenAIのDeepResearchの結果をプレゼンテーションに変換します。')
    parser.add_argument('input_file', help='入力テキストファイル（DeepResearchの結果）')
    parser.add_argument('-o', '--output', help='出力するパワーポイントファイル名')
    parser.add_argument('-t', '--title', help='プレゼンテーションのタイトル')
    parser.add_argument('--theme', choices=['blue', 'dark', 'light', 'green'], default='blue',
                        help='カラーテーマ（blue, dark, light, green）')
    
    args = parser.parse_args()
    
    try:
        with open(args.input_file, 'r', encoding='utf-8') as f:
            research_text = f.read()
        
        # タイトルが指定されていない場合はテキストから抽出
        title = args.title if args.title else extract_title_from_text(research_text)
        
        output_file = create_deep_research_presentation(
            research_text, 
            output_file=args.output,
            title=title,
            theme=args.theme
        )
        
        print(f"プレゼンテーションが正常に作成されました: {output_file}")
        
    except FileNotFoundError:
        print(f"エラー: ファイル '{args.input_file}' が見つかりません。")
    except Exception as e:
        print(f"エラー: {e}")

if __name__ == "__main__":
    main() 