from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import requests
from io import BytesIO

def create_arnak_presentation(output_file="arnak_strategy.pptx"):
    """
    アルナック（Lost Ruins of Arnak）の戦略をパワーポイントにまとめる関数
    
    Args:
        output_file (str): 出力するパワーポイントファイル名
    """
    # プレゼンテーションの作成
    prs = Presentation()
    
    # スライドのレイアウト
    title_slide_layout = prs.slide_layouts[0]  # タイトルスライド
    content_slide_layout = prs.slide_layouts[1]  # タイトルと内容のスライド
    
    # カラーテーマの設定（アルナックのイメージカラー）
    title_color = RGBColor(205, 133, 63)  # ペルー（明るいブラウン）
    subtitle_color = RGBColor(255, 215, 0)  # ゴールド
    text_color = RGBColor(255, 248, 220)  # コーンシルク（明るいベージュ）
    highlight_color = RGBColor(178, 34, 34)  # 赤茶色
    background_color = RGBColor(50, 25, 0)  # 暗いブラウン
    
    # 背景画像のダウンロードと保存
    bg_image_path = "arnak_bg.jpg"
    if not os.path.exists(bg_image_path):
        try:
            # 古代遺跡のイメージ画像をダウンロード
            bg_url = "https://images.unsplash.com/photo-1518998053901-5348d3961a04?q=80&w=1974&auto=format&fit=crop"
            response = requests.get(bg_url)
            with open(bg_image_path, "wb") as f:
                f.write(response.content)
            print(f"背景画像を {bg_image_path} として保存しました。")
        except Exception as e:
            print(f"背景画像のダウンロードに失敗しました: {e}")
            bg_image_path = None
    
    # タイトルスライドの作成
    title_slide = prs.slides.add_slide(title_slide_layout)
    
    # 背景画像の設定（タイトルスライド）
    if os.path.exists(bg_image_path):
        title_slide.shapes.add_picture(bg_image_path, 0, 0, prs.slide_width, prs.slide_height)
    
    # 半透明の背景オーバーレイを追加（タイトルが見やすくなるように）
    overlay = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill = overlay.fill
    fill.solid()
    fill.fore_color.rgb = background_color
    overlay.fill.transparency = 0.5
    
    # タイトルとサブタイトルのテキストボックスを追加
    title_box = title_slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(8), Inches(1.5)
    )
    title_tf = title_box.text_frame
    title_p = title_tf.add_paragraph()
    title_p.text = "アルナック"
    title_p.alignment = PP_ALIGN.CENTER
    title_run = title_p.runs[0]
    title_run.font.size = Pt(54)  # サイズを少し小さく
    title_run.font.bold = True
    title_run.font.color.rgb = title_color
    
    subtitle_box = title_slide.shapes.add_textbox(
        Inches(1), Inches(3), Inches(8), Inches(1)
    )
    subtitle_tf = subtitle_box.text_frame
    subtitle_p = subtitle_tf.add_paragraph()
    subtitle_p.text = "Lost Ruins of Arnak"
    subtitle_p.alignment = PP_ALIGN.CENTER
    subtitle_run = subtitle_p.runs[0]
    subtitle_run.font.size = Pt(32)  # サイズを少し小さく
    subtitle_run.font.italic = True
    subtitle_run.font.color.rgb = subtitle_color
    
    tagline_box = title_slide.shapes.add_textbox(
        Inches(1), Inches(4.5), Inches(8), Inches(1)
    )
    tagline_tf = tagline_box.text_frame
    tagline_p = tagline_tf.add_paragraph()
    tagline_p.text = "戦略ガイド"
    tagline_p.alignment = PP_ALIGN.CENTER
    tagline_run = tagline_p.runs[0]
    tagline_run.font.size = Pt(28)  # サイズを少し小さく
    tagline_run.font.bold = True
    tagline_run.font.color.rgb = text_color
    
    # 目次スライドの作成
    toc_slide = prs.slides.add_slide(content_slide_layout)
    
    # 背景の設定（目次スライド）
    if os.path.exists(bg_image_path):
        toc_slide.shapes.add_picture(bg_image_path, 0, 0, prs.slide_width, prs.slide_height)
    
    # 半透明の背景オーバーレイを追加
    overlay = toc_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill = overlay.fill
    fill.solid()
    fill.fore_color.rgb = background_color
    overlay.fill.transparency = 0.7
    
    # 目次タイトルのテキストボックス
    toc_title_box = toc_slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)  # 上部に配置
    )
    toc_title_tf = toc_title_box.text_frame
    toc_title_p = toc_title_tf.add_paragraph()
    toc_title_p.text = "目次"
    toc_title_p.alignment = PP_ALIGN.CENTER
    toc_title_run = toc_title_p.runs[0]
    toc_title_run.font.size = Pt(40)  # サイズを少し小さく
    toc_title_run.font.bold = True
    toc_title_run.font.color.rgb = title_color
    
    # 目次の内容テキストボックス
    toc_content_box = toc_slide.shapes.add_textbox(
        Inches(2), Inches(1.3), Inches(6), Inches(5.5)  # 縦幅を拡大
    )
    toc_content_tf = toc_content_box.text_frame
    
    # 目次項目
    toc_items = [
        "基本的な考え方",
        "序盤戦略",
        "中盤戦略",
        "終盤戦略",
        "上級者向けのポイント",
        "プレイヤー間のインタラクション",
        "注意点"
    ]
    
    for item in toc_items:
        p = toc_content_tf.add_paragraph()
        p.text = f"• {item}"
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(28)  # サイズを少し小さく
        run.font.color.rgb = text_color
        p.space_after = Pt(15)  # 間隔を少し狭く
    
    # コンテンツスライドの作成関数
    def create_content_slide(title, content_text):
        slide = prs.slides.add_slide(content_slide_layout)
        
        # 背景の設定
        if os.path.exists(bg_image_path):
            slide.shapes.add_picture(bg_image_path, 0, 0, prs.slide_width, prs.slide_height)
        
        # 半透明の背景オーバーレイを追加
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        fill = overlay.fill
        fill.solid()
        fill.fore_color.rgb = background_color
        overlay.fill.transparency = 0.7
        
        # タイトルのテキストボックス
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)  # 上部に配置
        )
        title_tf = title_box.text_frame
        title_p = title_tf.add_paragraph()
        title_p.text = title
        title_p.alignment = PP_ALIGN.CENTER
        title_run = title_p.runs[0]
        title_run.font.size = Pt(40)  # サイズを少し小さく
        title_run.font.bold = True
        title_run.font.color.rgb = title_color
        
        # 内容のテキストボックス
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2), Inches(9), Inches(5.5)  # 上部に配置し、縦幅を拡大
        )
        content_tf = content_box.text_frame
        content_tf.word_wrap = True
        
        # 内容のテキスト処理
        lines = content_text.split('\n')
        for i, line in enumerate(lines):
            # 空行はスキップ
            if not line.strip():
                continue
                
            p = content_tf.add_paragraph()
            p.text = line
            
            # インデントレベルに応じたスタイル設定
            if line.startswith('  - '):
                p.level = 1
                p.space_before = Pt(3)  # 間隔を狭く
                p.space_after = Pt(3)   # 間隔を狭く
            else:
                p.space_before = Pt(6)  # 間隔を狭く
                p.space_after = Pt(6)   # 間隔を狭く
            
            # runs[0]にアクセスする前にチェック
            if len(p.runs) > 0:
                run = p.runs[0]
                
                # 主要ポイントは強調
                if line.startswith('• '):
                    run.font.size = Pt(24)  # サイズを小さく
                    run.font.bold = True
                    run.font.color.rgb = text_color
                # サブポイントは少し小さく
                elif line.startswith('  - '):
                    run.font.size = Pt(20)  # サイズを小さく
                    run.font.color.rgb = text_color
                else:
                    run.font.size = Pt(24)  # サイズを小さく
                    run.font.color.rgb = text_color
        
        return slide
    
    # 基本的な考え方のスライド
    basic_content = """• ゲームは「研究トラックの先行」が勝利のカギ
• デッキ構築、ワーカープレイスメント、リソース管理の要素が融合している"""
    create_content_slide("基本的な考え方", basic_content)
    
    # 序盤戦略のスライド
    early_content = """• 研究トラックと助手の確保:
  - 序盤で研究トラックを進め、2ラウンド以内に助手2人を獲得
  - 助手は後半のアクション数や資源生産に大きく貢献

• アイテムカードの購入:
  - 強力なアイテム（犬、ランプ、双眼鏡、テントなど）を早めに入手
  - 資源獲得力を高める

• ワーカーの有効活用:
  - 遺跡発見や守護者討伐を目指し、初期段階で遺跡探索に着手
  - 偶像や守護者ボーナスを狙う"""
    create_content_slide("序盤戦略", early_content)
    
    # 中盤戦略のスライド
    mid_content = """• 研究トラックの継続:
  - 研究を進めて先行ボーナスを確保
  - 助手のアップグレードも進める
  - 後半の展開に直結

• 遺跡探索と守護者討伐:
  - レベル1・レベル2の遺跡発見・発掘をバランスよく行う
  - 高得点（偶像、守護者ボーナス）を狙う
  - 適切なリソース（コンパス、移動アイコンなど）の確保が重要

• デッキ圧縮:
  - 不要なカード（特に恐怖カードなど）を除去
  - 効率の良いデッキサイクルを維持"""
    create_content_slide("中盤戦略", mid_content)
    
    # 終盤戦略のスライド
    late_content = """• 研究トラックの仕上げ:
  - 虫眼鏡や手帳を頂上に到達させ、寺院タイルを獲得
  - 大幅な得点アップが期待できる

• リソースの使い切り:
  - 手元の余剰リソース（コイン、コンパス）はそのラウンド内にすべて活用
  - 追加の遺跡探索やアイテム購入、遺物活用へ変換

• 守護者と偶像:
  - 未討伐の守護者は最終ラウンドで確実に討伐
  - 偶像はそのまま得点化するか、状況に応じた追加効果に活用"""
    create_content_slide("終盤戦略", late_content)
    
    # 上級者向けのポイントのスライド
    advanced_content = """• 柔軟な判断:
  - 毎手番ごとに状況を分析し、最適なアクションを選択
  - 固定戦略に固執せず、相手の動向や市場状況を見極める

• リソースの相対評価:
  - リソース（コイン、コンパス、石版、矢じり、宝石）の価値を状況に応じて判断
  - 必要なものに集中的に投資

• カードドローとデッキ管理:
  - ドロー効果や除去効果をうまく活用
  - デッキの質を高めながら効率的な手番を実現"""
    create_content_slide("上級者向けのポイント", advanced_content)
    
    # プレイヤー間のインタラクションのスライド
    interaction_content = """• 研究トラックの競争:
  - 他プレイヤーよりも先に重要なマスに到達する
  - 相手の虫眼鏡/手帳コマの位置と所持リソースを常に確認

• ワーカープレイスメントのブロッキング:
  - 人気のキャンプ地や強力な遺跡は争奪戦になる
  - スタートプレイヤーの利点を活かす

• 市場のカード争奪:
  - 強力なアイテムは早い者勝ち
  - 「次のラウンドまで残らない」と思うカードは即座に確保"""
    create_content_slide("プレイヤー間のインタラクション", interaction_content)
    
    # 注意点のスライド
    caution_slide = prs.slides.add_slide(content_slide_layout)
    
    # 背景の設定
    if os.path.exists(bg_image_path):
        caution_slide.shapes.add_picture(bg_image_path, 0, 0, prs.slide_width, prs.slide_height)
    
    # 半透明の背景オーバーレイを追加
    overlay = caution_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill = overlay.fill
    fill.solid()
    fill.fore_color.rgb = background_color
    overlay.fill.transparency = 0.7
    
    # タイトルのテキストボックス
    caution_title_box = caution_slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)  # 上部に配置
    )
    caution_title_tf = caution_title_box.text_frame
    caution_title_p = caution_title_tf.add_paragraph()
    caution_title_p.text = "注意点"
    caution_title_p.alignment = PP_ALIGN.CENTER
    caution_title_run = caution_title_p.runs[0]
    caution_title_run.font.size = Pt(40)  # サイズを少し小さく
    caution_title_run.font.bold = True
    caution_title_run.font.color.rgb = highlight_color
    
    # 内容のテキストボックス
    caution_content_box = caution_slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(9), Inches(5.5)  # 上部に配置し、縦幅を拡大
    )
    caution_content_tf = caution_content_box.text_frame
    
    # 注意点の内容
    caution_points = [
        "助手確保の遅れ、研究トラックの軽視は致命的",
        "リソースの使い残しや無駄遣い、衝動的なカード購入は戦略全体を崩すリスクがある",
        "相手プレイヤーの動向を常に観察し、先手を取る意識が重要"
    ]
    
    for point in caution_points:
        p = caution_content_tf.add_paragraph()
        p.text = f"• {point}"
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(10)  # 間隔を狭く
        p.space_after = Pt(10)   # 間隔を狭く
        run = p.runs[0]
        run.font.size = Pt(24)   # サイズを小さく
        run.font.bold = True
        run.font.color.rgb = text_color
    
    # まとめスライドの作成
    summary_slide = prs.slides.add_slide(content_slide_layout)
    
    # 背景の設定
    if os.path.exists(bg_image_path):
        summary_slide.shapes.add_picture(bg_image_path, 0, 0, prs.slide_width, prs.slide_height)
    
    # 半透明の背景オーバーレイを追加
    overlay = summary_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill = overlay.fill
    fill.solid()
    fill.fore_color.rgb = background_color
    overlay.fill.transparency = 0.7
    
    # タイトルのテキストボックス
    summary_title_box = summary_slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)  # 上部に配置
    )
    summary_title_tf = summary_title_box.text_frame
    summary_title_p = summary_title_tf.add_paragraph()
    summary_title_p.text = "まとめ"
    summary_title_p.alignment = PP_ALIGN.CENTER
    summary_title_run = summary_title_p.runs[0]
    summary_title_run.font.size = Pt(40)  # サイズを少し小さく
    summary_title_run.font.bold = True
    summary_title_run.font.color.rgb = title_color
    
    # 内容のテキストボックス
    summary_content_box = summary_slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(9), Inches(5.5)  # 上部に配置し、縦幅を拡大
    )
    summary_content_tf = summary_content_box.text_frame
    
    # まとめの見出し
    summary_heading = summary_content_tf.add_paragraph()
    summary_heading.text = "アルナックの勝利の鍵："
    summary_heading.alignment = PP_ALIGN.CENTER
    summary_heading_run = summary_heading.runs[0]
    summary_heading_run.font.size = Pt(28)  # サイズを小さく
    summary_heading_run.font.bold = True
    summary_heading_run.font.color.rgb = subtitle_color
    summary_heading.space_after = Pt(15)  # 間隔を狭く
    
    # まとめの要点
    summary_points = [
        "序盤：研究トラックの先行と助手の確保",
        "中盤：研究継続と遺跡探索のバランス",
        "終盤：リソースの効率的な使い切りと得点の最大化",
        "常に：柔軟な判断と相手の動向観察"
    ]
    
    for point in summary_points:
        p = summary_content_tf.add_paragraph()
        p.text = f"• {point}"
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(8)   # 間隔を狭く
        p.space_after = Pt(8)    # 間隔を狭く
        run = p.runs[0]
        run.font.size = Pt(24)   # サイズを小さく
        run.font.bold = True
        run.font.color.rgb = text_color
    
    # プレゼンテーションの保存
    prs.save(output_file)
    print(f"アルナック戦略プレゼンテーションを {output_file} として保存しました。")

# メイン処理
if __name__ == "__main__":
    create_arnak_presentation() 