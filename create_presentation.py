from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re
import os

def create_board_game_presentation(game_info, output_file="board_game_strategy.pptx"):
    """
    ボードゲームの攻略情報をパワーポイントにまとめる関数
    
    Args:
        game_info (str): ボードゲームの攻略情報のテキスト
        output_file (str): 出力するパワーポイントファイル名
    """
    # プレゼンテーションの作成
    prs = Presentation()
    
    # スライドのレイアウト
    title_slide_layout = prs.slide_layouts[0]  # タイトルスライド
    content_slide_layout = prs.slide_layouts[1]  # タイトルと内容のスライド
    
    # タイトルスライドの作成
    title_slide = prs.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    # ゲーム名を抽出（最初の行または「ゲーム名:」などの形式から）
    game_name = game_info.strip().split('\n')[0]
    if ':' in game_name:
        game_name = game_name.split(':', 1)[1].strip()
    
    title.text = game_name
    subtitle.text = "ボードゲーム攻略ガイド"
    
    # テキストを段落に分割
    paragraphs = re.split(r'\n\s*\n', game_info)
    
    # 目次スライドの作成
    toc_slide = prs.slides.add_slide(content_slide_layout)
    toc_title = toc_slide.shapes.title
    toc_content = toc_slide.placeholders[1]
    
    toc_title.text = "目次"
    
    # 目次の内容を作成
    toc_text = ""
    section_titles = []
    
    for i, para in enumerate(paragraphs):
        if i == 0:  # 最初の段落はタイトルなのでスキップ
            continue
        
        # 段落の最初の行をセクションタイトルとして使用
        lines = para.strip().split('\n')
        if lines:
            section_title = lines[0].strip()
            # URLを含む行は除外
            if not section_title.startswith('http') and not 'www.' in section_title:
                section_titles.append(section_title)
                toc_text += f"• {section_title}\n"
    
    toc_content.text = toc_text
    
    # 各セクションのスライドを作成
    current_section = ""
    section_content = ""
    
    for i, para in enumerate(paragraphs):
        if i == 0:  # 最初の段落はタイトルなのでスキップ
            continue
        
        lines = para.strip().split('\n')
        if not lines:
            continue
        
        # URLを含む行は除外
        filtered_lines = [line for line in lines if not line.startswith('http') and not 'www.' in line]
        if not filtered_lines:
            continue
        
        section_title = filtered_lines[0].strip()
        
        # 新しいセクションの開始
        if section_title in section_titles:
            # 前のセクションがあれば、そのスライドを作成
            if current_section:
                create_section_slide(prs, content_slide_layout, current_section, section_content)
            
            current_section = section_title
            section_content = '\n'.join(filtered_lines[1:])
        else:
            # 同じセクションの続き
            section_content += '\n' + '\n'.join(filtered_lines)
    
    # 最後のセクションのスライドを作成
    if current_section:
        create_section_slide(prs, content_slide_layout, current_section, section_content)
    
    # まとめスライドの作成
    summary_slide = prs.slides.add_slide(content_slide_layout)
    summary_title = summary_slide.shapes.title
    summary_content = summary_slide.placeholders[1]
    
    summary_title.text = "まとめ"
    summary_content.text = f"{game_name}の攻略ポイント：\n\n• 基本ルールを理解する\n• 戦略的な思考を身につける\n• 経験を積んで上達しよう"
    
    # プレゼンテーションの保存
    prs.save(output_file)
    print(f"プレゼンテーションを {output_file} として保存しました。")

def create_section_slide(prs, layout, title, content):
    """
    セクションのスライドを作成する関数
    
    Args:
        prs: プレゼンテーションオブジェクト
        layout: スライドレイアウト
        title (str): スライドのタイトル
        content (str): スライドの内容
    """
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    
    # 内容が長すぎる場合は分割
    if len(content) > 1000:
        # 内容を箇条書きに変換
        bullet_points = []
        for line in content.split('\n'):
            line = line.strip()
            if line:
                if not line.startswith('•'):
                    line = f"• {line}"
                bullet_points.append(line)
        
        # 箇条書きを結合
        content = '\n'.join(bullet_points[:10])  # 最初の10項目だけ表示
        content += "\n• ..."  # 省略記号を追加
    
    content_shape.text = content

def create_research_presentation(research_text, output_file="research_presentation.pptx", title="研究結果"):
    """
    OpenAIのDeepResearchの結果をパワーポイントにまとめる関数
    
    Args:
        research_text (str): DeepResearchの結果テキスト
        output_file (str): 出力するパワーポイントファイル名
        title (str): プレゼンテーションのタイトル
    """
    # プレゼンテーションの作成
    prs = Presentation()
    
    # スライドのレイアウト
    title_slide_layout = prs.slide_layouts[0]  # タイトルスライド
    content_slide_layout = prs.slide_layouts[1]  # タイトルと内容のスライド
    
    # カラーテーマの設定
    title_color = RGBColor(0, 112, 192)  # 青
    subtitle_color = RGBColor(0, 176, 240)  # 明るい青
    text_color = RGBColor(0, 0, 0)  # 黒
    highlight_color = RGBColor(192, 0, 0)  # 赤
    
    # タイトルスライドの作成
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_shape = title_slide.shapes.title
    subtitle_shape = title_slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = "研究結果プレゼンテーション"
    
    # URLを除外するための正規表現パターン
    url_pattern = re.compile(r'https?://\S+|www\.\S+|\[\d+\]|\(\d+\)|参考文献|References')
    
    # テキストを段落に分割
    paragraphs = re.split(r'\n\s*\n', research_text)
    
    # 目次スライドの作成
    toc_slide = prs.slides.add_slide(content_slide_layout)
    toc_title = toc_slide.shapes.title
    toc_content = toc_slide.placeholders[1]
    
    toc_title.text = "目次"
    
    # 目次の内容を作成
    toc_text = ""
    section_titles = []
    
    for i, para in enumerate(paragraphs):
        if i == 0:  # 最初の段落はタイトルなのでスキップ
            continue
        
        # 段落の最初の行をセクションタイトルとして使用
        lines = para.strip().split('\n')
        if lines:
            section_title = lines[0].strip()
            # URLや参考文献を含む行は除外
            if not url_pattern.search(section_title):
                # 数字だけの見出しや短すぎる見出しは除外
                if not re.match(r'^\d+\.?$', section_title) and len(section_title) > 3:
                    section_titles.append(section_title)
                    toc_text += f"• {section_title}\n"
    
    toc_content.text = toc_text
    
    # 各セクションのスライドを作成
    current_section = ""
    section_content = ""
    
    for i, para in enumerate(paragraphs):
        if i == 0:  # 最初の段落はタイトルなのでスキップ
            continue
        
        lines = para.strip().split('\n')
        if not lines:
            continue
        
        # URLや参考文献を含む行は除外
        filtered_lines = []
        for line in lines:
            # 参考文献セクションを検出したら、そのパラグラフ以降は処理しない
            if re.match(r'^参考文献|^References', line, re.IGNORECASE):
                break
            # URLや引用番号を含まない行だけを追加
            if not url_pattern.search(line):
                filtered_lines.append(line)
        
        if not filtered_lines:
            continue
        
        section_title = filtered_lines[0].strip()
        
        # 新しいセクションの開始
        if section_title in section_titles:
            # 前のセクションがあれば、そのスライドを作成
            if current_section:
                create_research_slide(prs, content_slide_layout, current_section, section_content, title_color, text_color)
            
            current_section = section_title
            section_content = '\n'.join(filtered_lines[1:])
        else:
            # 同じセクションの続き
            section_content += '\n' + '\n'.join(filtered_lines)
    
    # 最後のセクションのスライドを作成
    if current_section:
        create_research_slide(prs, content_slide_layout, current_section, section_content, title_color, text_color)
    
    # まとめスライドの作成
    summary_slide = prs.slides.add_slide(content_slide_layout)
    summary_title = summary_slide.shapes.title
    summary_content = summary_slide.placeholders[1]
    
    summary_title.text = "まとめ"
    
    # まとめの内容（最初の段落から抽出するか、固定テキスト）
    if len(paragraphs) > 0:
        # 最初の段落から要約を抽出（URLや参考文献を除外）
        first_para = paragraphs[0]
        summary_lines = [line for line in first_para.split('\n') if not url_pattern.search(line)]
        summary_text = '\n'.join(summary_lines)
        
        # 長すぎる場合は短縮
        if len(summary_text) > 500:
            summary_text = summary_text[:497] + "..."
    else:
        summary_text = "研究結果の主要ポイント"
    
    summary_content.text = summary_text
    
    # プレゼンテーションの保存
    prs.save(output_file)
    print(f"研究プレゼンテーションを {output_file} として保存しました。")

def create_research_slide(prs, layout, title, content, title_color, text_color):
    """
    研究結果のスライドを作成する関数
    
    Args:
        prs: プレゼンテーションオブジェクト
        layout: スライドレイアウト
        title (str): スライドのタイトル
        content (str): スライドの内容
        title_color: タイトルの色
        text_color: テキストの色
    """
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    
    # タイトルの色を設定
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = title_color
    
    # 内容が長すぎる場合は分割
    if len(content) > 1000:
        # 内容を箇条書きに変換
        bullet_points = []
        for line in content.split('\n'):
            line = line.strip()
            if line:
                # すでに箇条書きになっている場合はそのまま、そうでなければ箇条書きに変換
                if not line.startswith('•') and not line.startswith('-'):
                    line = f"• {line}"
                bullet_points.append(line)
        
        # 箇条書きを結合（最初の部分だけ表示）
        content = '\n'.join(bullet_points[:15])  # 最初の15項目だけ表示
        if len(bullet_points) > 15:
            content += "\n• ..."  # 省略記号を追加
    
    content_shape.text = content
    
    # テキストの色を設定
    for paragraph in content_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = text_color
    
    return slide

# メイン処理
if __name__ == "__main__":
    # ユーザーからの入力を受け取る
    print("ボードゲームの攻略情報またはDeepResearchの結果をテキストファイルから読み込みます。")
    
    try:
        with open("game_info.txt", "r", encoding="utf-8") as f:
            text_content = f.read()
        
        # 内容に基づいて適切な関数を呼び出す
        if "参考文献" in text_content or "References" in text_content or "http" in text_content:
            print("DeepResearchの結果と判断しました。研究プレゼンテーションを作成します。")
            create_research_presentation(text_content)
        else:
            print("ボードゲームの攻略情報と判断しました。ボードゲームプレゼンテーションを作成します。")
            create_board_game_presentation(text_content)
    except FileNotFoundError:
        print("game_info.txt ファイルが見つかりません。")
        print("テキストファイルを作成し、ボードゲームの攻略情報またはDeepResearchの結果を記入してください。") 